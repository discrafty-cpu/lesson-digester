#!/usr/bin/env python3
"""
THE LESSON DIGESTER — Drummond Design System Engine
Built with python-pptx

Chews up cluttered CPM lesson decks and spits out clean, standards-aligned,
story-driven presentations with embedded collaborative roles.

Usage:
    python engine.py              # Build demo deck
    python engine.py input.pptx   # Analyze + rebuild an existing deck (future)

Each build_* function creates one slide type. Customize the data dicts to match your lesson.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, sys

# ─── DESIGN TOKENS ──────────────────────────────────────────
class C:
    """Drummond Design System Colors"""
    navy       = RGBColor(0x1E, 0x27, 0x61)
    charcoal   = RGBColor(0x2D, 0x37, 0x48)
    amber      = RGBColor(0xD4, 0x87, 0x0F)
    teal       = RGBColor(0x0D, 0x94, 0x88)
    blue       = RGBColor(0x3B, 0x82, 0xF6)
    white      = RGBColor(0xFF, 0xFF, 0xFF)
    surface    = RGBColor(0xF7, 0xF8, 0xFA)
    green      = RGBColor(0x4C, 0xAF, 0x50)
    red        = RGBColor(0xEF, 0x44, 0x44)
    slate      = RGBColor(0x64, 0x74, 0x8B)
    darkBg     = RGBColor(0x1A, 0x20, 0x2C)
    greenPanel = RGBColor(0x86, 0xB0, 0x49)
    lightBlue  = RGBColor(0xEB, 0xF5, 0xFF)
    lightAmber = RGBColor(0xFF, 0xF8, 0xE1)
    cream      = RGBColor(0xF5, 0xF0, 0xE8)
    # Role colors
    facilitator = RGBColor(0x0D, 0x94, 0x88)  # teal
    resource_mgr = RGBColor(0xD4, 0x87, 0x0F)  # amber
    recorder    = RGBColor(0x3B, 0x82, 0xF6)   # blue
    task_mgr    = RGBColor(0xE8, 0x43, 0x6D)   # pink/red

FONT_TITLE = "Arial Black"
FONT_HEADING = "Calibri"
FONT_BODY = "Calibri"

# ─── HELPERS ─────────────────────────────────────────────────

def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text="", font_size=14,
                font_name=FONT_BODY, color=None, bold=False, italic=False,
                align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, word_wrap=True):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.color.rgb = color or C.charcoal
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = align
    tf.auto_size = None
    # Vertical alignment
    txBox.text_frame._txBody.attrib['{http://schemas.openxmlformats.org/drawingml/2006/main}anchor'] = {
        MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'
    }.get(valign, 't')
    return txBox

def add_rich_textbox(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT):
    """Add a text box with multiple formatted runs.
    runs: list of dicts with keys: text, font_size, font_name, color, bold, italic, breakLine
    """
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    current_para = tf.paragraphs[0]
    current_para.alignment = align

    for i, run_data in enumerate(runs):
        if i > 0 and run_data.get('breakLine', False):
            current_para = tf.add_paragraph()
            current_para.alignment = align

        run = current_para.add_run()
        run.text = run_data.get('text', '')
        run.font.size = Pt(run_data.get('font_size', 14))
        run.font.name = run_data.get('font_name', FONT_BODY)
        run.font.color.rgb = run_data.get('color', C.charcoal)
        run.font.bold = run_data.get('bold', False)
        run.font.italic = run_data.get('italic', False)

    return txBox

def add_line(slide, left, top, width):
    """Add a horizontal line."""
    connector = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(3)
    )
    connector.fill.solid()
    connector.fill.fore_color.rgb = C.amber
    connector.line.fill.background()
    return connector

def add_accent_bar(slide, left, top, height, color=None):
    """Add a vertical accent bar (Drummond signature element)."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.06), Inches(height)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color or C.teal
    bar.line.fill.background()
    return bar


# ─── SLIDE BUILDERS ──────────────────────────────────────────

def build_intro_agenda(prs, data):
    """INTRO_AGENDA: NEED / DO / NEXT header + Focus Skill"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, C.surface)

    # Top strip: NEED / DO / NEXT cards
    card_y = 0.2
    card_h = 1.3
    card_w = 2.9
    gap = 0.15

    cards = [
        ("NEED", data.get("need", "Notebook, Pencil, Chromebook")),
        ("DO", data.get("do_text", 'Label Your Notebook:\nUnit 4: 7.1.1\n"Distance Rate and Time"')),
        ("NEXT", data.get("next", "Focus Skill:\n1. Draw a line diagram\n2. Solve for the missing value")),
    ]

    for i, (label, body) in enumerate(cards):
        x = 0.3 + i * (card_w + gap + 0.3)
        # Card bg
        add_shape(slide, x, card_y, card_w, card_h, fill_color=C.white)
        # Label
        add_textbox(slide, x + 0.15, card_y + 0.08, card_w - 0.3, 0.35,
                    text=label, font_size=16, font_name=FONT_TITLE,
                    color=C.amber, bold=True)
        # Body
        add_textbox(slide, x + 0.15, card_y + 0.45, card_w - 0.3, card_h - 0.55,
                    text=body, font_size=11, color=C.charcoal)
        # Arrow between cards
        if i < 2:
            arrow = add_shape(slide, x + card_w + 0.02, card_y + card_h / 2 - 0.03,
                              0.26, 0.06, fill_color=C.amber)

    # Focus Skill section
    fs_y = 1.75
    add_shape(slide, 0.3, fs_y, 9.4, 3.6, fill_color=C.white)
    add_accent_bar(slide, 0.3, fs_y, 0.45, C.teal)

    # Focus Skill title
    add_rich_textbox(slide, 0.5, fs_y + 0.05, 8, 0.4, [
        {"text": "Focus Skill: ", "bold": True, "color": C.navy, "font_size": 22, "font_name": FONT_HEADING},
        {"text": data.get("focus_title", "Review"), "bold": True, "color": C.amber, "font_size": 22, "font_name": FONT_HEADING},
    ])

    # Focus problems (two-column)
    problems = data.get("focus_problems", [])
    for i, prob in enumerate(problems[:2]):
        px = 0.5 + i * 4.7
        add_textbox(slide, px, fs_y + 0.55, 4.2, 0.3,
                    text=prob.get("title", f"Problem {i+1}"),
                    font_size=14, font_name=FONT_HEADING, bold=True, color=C.navy)
        add_textbox(slide, px, fs_y + 0.9, 4.2, 2.2,
                    text=prob.get("text", ""),
                    font_size=12, color=C.charcoal)

    # Instructions
    instructions = data.get("focus_instructions", [])
    if instructions:
        inst_text = "\n".join(f"{i+1}. {inst}" for i, inst in enumerate(instructions))
        add_textbox(slide, 2.5, fs_y + 3.1, 5, 0.45,
                    text=inst_text, font_size=11, color=C.slate)

    return slide


def build_checkin(prs, data):
    """CHECKIN_FEELINGS: Social-emotional check-in"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.darkBg)

    add_textbox(slide, 0.5, 0.3, 9, 0.6,
                text="Share with your TEAM", font_size=32,
                font_name=FONT_TITLE, color=C.white, align=PP_ALIGN.CENTER)

    add_textbox(slide, 1, 1.1, 8, 0.5,
                text="What number/picture best shows how you are feeling?  Why?",
                font_size=16, color=C.slate, italic=True, align=PP_ALIGN.CENTER)

    # Sentence frame card
    add_shape(slide, 1.5, 1.9, 7, 0.7, fill_color=RGBColor(0x2D, 0x37, 0x48))
    add_textbox(slide, 1.7, 1.95, 6.6, 0.6,
                text='"I am feeling like number _____ because _____"',
                font_size=16, color=C.amber, italic=True, align=PP_ALIGN.CENTER)

    # Sharing order table
    from pptx.util import Inches as In, Pt as P
    rows, cols = 2, 2
    tbl = slide.shapes.add_table(rows, cols, In(2.5), In(3.0), In(5), In(1.6)).table
    role_data = [
        [("Recorder Reporter\n4th to share", "FFE0E0"), ("Facilitator\n1st to share", "E0F0E0")],
        [("Task Manager\n3rd to share", "E0E0FF"), ("Resource Manager\n2nd to share", "FFF0E0")],
    ]
    for r in range(rows):
        for c_ in range(cols):
            cell = tbl.cell(r, c_)
            text, bg = role_data[r][c_]
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(bg)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = C.charcoal
                p.alignment = PP_ALIGN.CENTER

    return slide


def build_breathing(prs, data):
    """BREATHING_MINDFULNESS: Box breathing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(0x1B, 0x28, 0x38))

    add_textbox(slide, 1, 0.5, 8, 0.6,
                text="Box Breathing", font_size=36,
                font_name=FONT_TITLE, color=C.white, align=PP_ALIGN.CENTER)

    steps = [
        ("Breathe IN for 4 counts", False),
        ("HOLD for 4 counts", False),
        ("Breathe OUT for 4 counts", False),
        ("HOLD for 4 counts", False),
        ("Repeat", True),
    ]
    runs = []
    for text, is_accent in steps:
        runs.append({
            "text": text, "breakLine": True,
            "font_size": 20, "color": C.amber if is_accent else C.white,
            "bold": is_accent,
        })
    add_rich_textbox(slide, 2.5, 1.5, 5, 3, runs, align=PP_ALIGN.CENTER)
    return slide


def build_learning_target(prs, data):
    """LEARNING_TARGET: Objectives + Success Criteria"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.surface)

    add_textbox(slide, 0.4, 0.2, 5.5, 0.5,
                text="Learning Target:", font_size=28,
                font_name=FONT_TITLE, color=C.navy)

    # "I will..." statements
    targets = data.get("targets", [])
    runs = []
    for i, t in enumerate(targets):
        if i > 0:
            runs.append({"text": "", "breakLine": True, "font_size": 6, "color": C.charcoal})
        runs.append({"text": "I will ", "font_size": 15, "color": C.charcoal, "breakLine": i > 0})
        runs.append({"text": t.get("verb", ""), "font_size": 15, "color": C.navy, "bold": True})
        runs.append({"text": f" {t.get('rest', '')}", "font_size": 15, "color": C.charcoal})
    if runs:
        add_rich_textbox(slide, 0.4, 0.8, 5.5, 1.8, runs)

    # Discussion prompt
    discuss = data.get("discuss_prompt")
    if discuss:
        add_shape(slide, 0.4, 2.8, 5.5, 1.4, fill_color=C.lightBlue)
        add_rich_textbox(slide, 0.6, 2.85, 5.1, 0.6, [
            {"text": "Discuss:\n", "bold": True, "font_size": 14, "color": C.navy},
            {"text": discuss, "font_size": 13, "color": C.charcoal},
        ])
        frame = data.get("discuss_frame")
        if frame:
            add_textbox(slide, 0.6, 3.6, 5.1, 0.4,
                        text=f'"{frame}"', font_size=12, color=C.teal, italic=True)

    # Success Criteria card (right)
    add_shape(slide, 6.2, 0.2, 3.5, 3.2, fill_color=C.white, line_color=C.blue, line_width=2)
    add_textbox(slide, 6.35, 0.3, 3.2, 0.35,
                text="Success Criteria:", font_size=16,
                font_name=FONT_HEADING, bold=True, color=C.navy)
    add_textbox(slide, 6.35, 0.65, 3.2, 0.25,
                text="I'll know if I've got it when:", font_size=10, bold=True, color=C.slate)

    criteria = data.get("success_criteria", [])
    crit_text = "\n".join(f"  {sc}" for sc in criteria)
    add_textbox(slide, 6.35, 1.0, 3.2, 2.2,
                text=crit_text, font_size=10, color=C.charcoal)

    return slide


def build_story_narrative(prs, data):
    """STORY_NARRATIVE: Story hook with discussion prompt"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.white)

    # Story title (teal/mint accent)
    add_textbox(slide, 0.5, 0.3, 6, 0.7,
                text=data.get("title", "The Story Begins"),
                font_size=36, font_name=FONT_TITLE,
                color=RGBColor(0x2D, 0xD4, 0xBF))

    # Story body
    add_textbox(slide, 0.5, 1.2, 5.5, 2.5,
                text=data.get("text", ""), font_size=16, color=C.charcoal)

    # Discussion box
    discuss = data.get("discuss")
    if discuss:
        add_shape(slide, 5.5, 3.5, 4.2, 1.8, fill_color=C.lightBlue)
        add_rich_textbox(slide, 5.7, 3.6, 3.8, 1.5, [
            {"text": "Discuss: ", "bold": True, "font_size": 13, "color": C.navy},
            {"text": discuss, "font_size": 13, "color": C.red},
        ])

    # Image
    img = data.get("image")
    if img and os.path.exists(img):
        slide.shapes.add_picture(img, Inches(6.3), Inches(0.5), Inches(3.2), Inches(2.8))

    return slide


def build_story_dramatic(prs, data):
    """STORY_DRAMATIC: Dark dramatic beat with alert"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.darkBg)

    # Left: dramatic title
    add_textbox(slide, 0.3, 0.3, 4.5, 0.7,
                text=data.get("title", "ALERT"),
                font_size=32, font_name=FONT_TITLE, color=C.red)
    add_textbox(slide, 0.3, 1.2, 4.5, 0.7,
                text=data.get("subtitle", ""),
                font_size=18, font_name=FONT_HEADING, bold=True, color=C.white)

    # Right: green panel with story continuation
    add_shape(slide, 5.0, 0, 5, 3.0, fill_color=C.greenPanel)
    add_textbox(slide, 5.2, 0.2, 4.6, 2.6,
                text=data.get("quote", ""),
                font_size=14, color=C.white)

    # Role prompt (above details, not overlapping)
    role_prompt = data.get("role_prompt")
    detail_start_y = 3.2
    if role_prompt:
        add_textbox(slide, 5.2, 3.2, 4.6, 0.4,
                    text=role_prompt, font_size=11,
                    color=C.amber, bold=True, italic=True)
        detail_start_y = 3.65

    # Details (below role prompt)
    details = data.get("details", [])
    if details:
        detail_text = "\n".join(f"  {d}" for d in details)
        add_textbox(slide, 5.2, detail_start_y, 4.6, 1.8,
                    text=detail_text, font_size=13,
                    color=RGBColor(0xCB, 0xD5, 0xE0))

    return slide


def build_story_mission(prs, data):
    """STORY_MISSION: Team discussion bridge (story to math) with embedded role prompts"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.white)

    # Left dark panel
    add_shape(slide, 0, 0, 5, 5.625, fill_color=C.darkBg)
    add_accent_bar(slide, 0.4, 1.0, 0.4, C.teal)
    add_textbox(slide, 0.6, 0.95, 4, 0.5,
                text=data.get("title", "The Mission"),
                font_size=28, font_name=FONT_TITLE, bold=True, color=C.white)
    add_textbox(slide, 0.6, 1.7, 4, 1.5,
                text=data.get("quote", ""), font_size=14,
                color=RGBColor(0xCB, 0xD5, 0xE0))

    # Task Manager callout box
    tm_prompt = data.get("task_manager_prompt")
    if tm_prompt:
        add_shape(slide, 0.4, 3.5, 4.2, 0.8, fill_color=RGBColor(0x2D, 0x37, 0x48))
        add_rich_textbox(slide, 0.55, 3.55, 3.9, 0.7, [
            {"text": "Task Manager: ", "bold": True, "italic": True, "font_size": 11, "color": C.task_mgr},
            {"text": tm_prompt, "italic": True, "font_size": 11, "color": RGBColor(0xCB, 0xD5, 0xE0)},
        ])

    # Right: role-specific discussion prompts
    role_prompts = data.get("role_prompts", [])
    role_colors = {
        "Facilitator": C.facilitator,
        "Resource Manager": C.resource_mgr,
        "Task Manager": C.task_mgr,
        "Recorder Reporter": C.recorder,
    }
    for i, rp in enumerate(role_prompts[:4]):
        ry = 0.3 + i * 1.1
        role_name = rp.get("role", "")
        role_color = role_colors.get(role_name, C.charcoal)
        add_rich_textbox(slide, 5.2, ry, 4.5, 0.9, [
            {"text": f"{role_name}: ", "bold": True, "italic": True, "font_size": 12, "color": role_color},
            {"text": f'"{rp.get("prompt", "")}"', "italic": True, "font_size": 12, "color": C.charcoal},
        ])

    return slide


def build_group_discussion(prs, data):
    """GROUP_DISCUSSION: Guided team conversation with role-specific prompt cards"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.surface)

    # Header with accent bar
    add_accent_bar(slide, 0.3, 0.2, 0.4, C.teal)
    add_textbox(slide, 0.5, 0.2, 9, 0.4,
                text=data.get("title", "Team Discussion"),
                font_size=22, font_name=FONT_HEADING, bold=True, color=C.navy)

    # Main question
    add_textbox(slide, 0.5, 0.7, 9, 0.5,
                text=data.get("question", ""), font_size=16, color=C.charcoal)

    # Role-specific prompt cards (2x2)
    role_cards = data.get("guided_prompts", [
        {"role": "Facilitator", "prompt": "Read the question aloud. Ask: 'What do we think?'", "color": C.facilitator},
        {"role": "Resource Manager", "prompt": "What tools or information do we need?", "color": C.resource_mgr},
        {"role": "Task Manager", "prompt": "Make sure everyone shares. Watch the time.", "color": C.task_mgr},
        {"role": "Recorder Reporter", "prompt": "Write down the team's answer.", "color": C.recorder},
    ])

    for i, rc in enumerate(role_cards[:4]):
        cx = 0.3 + (i % 2) * 4.85
        cy = 1.5 + (i // 2) * 1.85
        color = rc.get("color", C.charcoal)
        if isinstance(color, str):
            color = RGBColor.from_string(color)

        # Card
        add_shape(slide, cx, cy, 4.55, 1.6, fill_color=C.white)
        # Color accent at top
        add_shape(slide, cx, cy, 4.55, 0.05, fill_color=color)
        # Role name
        add_textbox(slide, cx + 0.15, cy + 0.12, 4.2, 0.3,
                    text=rc["role"], font_size=14, font_name=FONT_HEADING,
                    bold=True, color=color)
        # Prompt
        add_textbox(slide, cx + 0.15, cy + 0.5, 4.2, 0.9,
                    text=f'"{rc["prompt"]}"', font_size=12,
                    color=C.charcoal, italic=True)

    # Sentence frame at bottom
    frame = data.get("sentence_frame")
    if frame:
        add_shape(slide, 0.3, 5.0, 9.4, 0.45, fill_color=C.lightAmber)
        add_textbox(slide, 0.5, 5.02, 9, 0.4,
                    text=f'"{frame}"', font_size=12, color=C.amber,
                    italic=True, align=PP_ALIGN.CENTER)

    return slide


def build_activity_launch(prs, data):
    """ACTIVITY_LAUNCH: DO slide with role assignments"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.surface)

    add_textbox(slide, 0.3, 0.2, 1.2, 0.5,
                text="DO:", font_size=28, font_name=FONT_TITLE, color=C.navy)

    # Teams Complete card
    add_shape(slide, 0.3, 0.8, 5.5, 1.8, fill_color=C.white)
    add_textbox(slide, 2.0, 0.85, 1.5, 0.4,
                text="CPM", font_size=18, font_name=FONT_TITLE,
                color=C.teal, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 1.3, 5, 0.8,
                text=f"Teams Complete:\n{data.get('ref', 'Unit - Lesson: Problems')}",
                font_size=13, bold=True, color=C.charcoal)

    # Role assignment table
    roles = data.get("roles", [["Computers", "Read 1st Q"], ["Calculators", "Responds 1st"]])
    tbl = slide.shapes.add_table(len(roles), 2, Inches(0.5), Inches(2.2), Inches(4), Inches(0.8)).table
    for r, row in enumerate(roles):
        for c_, val in enumerate(row):
            cell = tbl.cell(r, c_)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = C.charcoal
                p.alignment = PP_ALIGN.CENTER

    # Need section
    add_textbox(slide, 6.2, 0.3, 3.5, 0.35,
                text="Need:", font_size=16, font_name=FONT_HEADING,
                bold=True, color=C.navy)
    needs = data.get("needs", ["Chromebook", "Calculator", "Notebook"])
    add_textbox(slide, 6.2, 0.7, 3.5, 2.0,
                text="\n".join(f"  {n}" for n in needs),
                font_size=12, color=C.charcoal)

    return slide


def build_problem(prs, data):
    """PROBLEM_SLIDE: Individual problem with embedded role prompts"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.white)

    # Problem number header
    add_textbox(slide, 0.3, 0.2, 9.4, 0.6,
                text=data.get("number", "7-1"),
                font_size=28, font_name=FONT_TITLE, color=C.navy)
    add_line(slide, 0.3, 0.85, 2)

    # Problem text
    add_textbox(slide, 0.5, 1.1, 6.5, 3.0,
                text=data.get("text", ""), font_size=16, color=C.charcoal)

    # Embedded role prompt strip (bottom) - contextual, not boilerplate
    role_strip = data.get("role_prompts")
    if role_strip:
        strip_y = 4.3
        add_shape(slide, 0.3, strip_y, 9.4, 1.1, fill_color=C.surface)
        role_colors = [C.facilitator, C.resource_mgr, C.task_mgr, C.recorder]
        for i, (role, prompt) in enumerate(role_strip[:4]):
            rx = 0.4 + i * 2.35
            # Color dot
            add_shape(slide, rx, strip_y + 0.1, 0.12, 0.12, fill_color=role_colors[i])
            add_textbox(slide, rx + 0.18, strip_y + 0.05, 2.0, 0.2,
                        text=role, font_size=8, bold=True, color=role_colors[i])
            add_textbox(slide, rx + 0.18, strip_y + 0.25, 2.0, 0.7,
                        text=prompt, font_size=8, color=C.slate, italic=True)

    # Image
    img = data.get("image")
    if img and os.path.exists(img):
        slide.shapes.add_picture(img, Inches(7.0), Inches(1.0), Inches(2.7), Inches(2.5))

    return slide


def build_closure(prs, data):
    """CLOSURE: End-of-lesson discussion"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.surface)

    add_textbox(slide, 0.3, 0.2, 9.4, 0.7,
                text="CLOSURE", font_size=36, font_name=FONT_TITLE, color=C.navy)
    add_line(slide, 0.3, 0.95, 3)

    # Discussion card
    add_shape(slide, 0.5, 1.3, 9, 2.5, fill_color=C.white)
    add_textbox(slide, 0.8, 1.5, 8.4, 2.0,
                text=data.get("prompt", "Discuss: What did you learn today?"),
                font_size=20, color=C.charcoal)

    return slide


def build_homework(prs, data):
    """HOMEWORK: Tonight's assignment"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.navy)

    # Banner
    add_shape(slide, 1.5, 0.8, 7, 1.2, fill_color=C.amber)
    add_textbox(slide, 1.5, 0.85, 7, 1.1,
                text="HOMEWORK TONIGHT", font_size=36,
                font_name=FONT_TITLE, color=C.white,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # Assignment
    add_textbox(slide, 1.5, 2.5, 7, 1.0,
                text=data.get("ref", "Lesson: Problems"),
                font_size=28, font_name=FONT_HEADING,
                color=C.white, bold=True, align=PP_ALIGN.CENTER)

    return slide


def build_day_divider(prs, data):
    """DAY_DIVIDER: Day 2 separator"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.navy)

    add_textbox(slide, 0, 1.0, 10, 2.5,
                text=data.get("label", "Day 2"),
                font_size=72, font_name=FONT_TITLE,
                color=C.white, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_line(slide, 4, 3.7, 2)
    add_textbox(slide, 1, 3.9, 8, 0.6,
                text=data.get("subtitle", ""),
                font_size=18, color=C.slate, align=PP_ALIGN.CENTER)

    return slide


def build_team_roles(prs, data):
    """TEAM_ROLES: 2x2 role cards on dark background (only for new-unit/new-seats)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.darkBg)

    roles = data.get("roles", [
        {"name": "Recorder & Reporter", "name_es": "Registrador y Relator", "color": C.recorder,
         "jobs": ["Record all final answers and methods.", "Share the team's ideas with the class.", "Make sure the team's work is clear."],
         "jobs_es": ["Registrar todas las respuestas.", "Compartir las ideas del equipo.", "Asegurar claridad."]},
        {"name": "Resource Manager", "name_es": "Gerente de Recursos", "color": C.resource_mgr,
         "jobs": ["Get supplies for the team.", "Call the teacher for questions.", "Make sure everything is cleaned up."],
         "jobs_es": ["Obtener suministros.", "Llamar al maestro.", "Asegurar limpieza."]},
        {"name": "Task Manager", "name_es": "Gerente de Tareas", "color": C.task_mgr,
         "jobs": ["Keep the team focused on the math.", "Watch the clock and time limits.", "Read instructions out loud."],
         "jobs_es": ["Mantener enfoque.", "Vigilar el reloj.", "Leer instrucciones."]},
        {"name": "Facilitator", "name_es": "Facilitador", "color": C.facilitator,
         "jobs": ["Start by reading the problem.", "Make sure everyone can explain.", '"Does everyone understand?"'],
         "jobs_es": ["Iniciar leyendo el problema.", "Asegurar comprension.", '"Todos entienden?"']},
    ])

    positions = [(0.3, 0.2), (5.15, 0.2), (0.3, 2.9), (5.15, 2.9)]
    card_w, card_h = 4.55, 2.5

    for i, role in enumerate(roles[:4]):
        x, y = positions[i]
        color = role["color"]
        if isinstance(color, str):
            color = RGBColor.from_string(color)

        # Card bg
        add_shape(slide, x, y, card_w, card_h, fill_color=RGBColor(0x2D, 0x37, 0x48))
        # Color accent top
        add_shape(slide, x, y, card_w, 0.06, fill_color=color)
        # Role name
        add_textbox(slide, x + 0.15, y + 0.15, card_w - 0.3, 0.3,
                    text=role["name"], font_size=15, font_name=FONT_HEADING,
                    bold=True, color=C.white)
        add_textbox(slide, x + 0.15, y + 0.45, card_w - 0.3, 0.2,
                    text=role.get("name_es", ""), font_size=10, color=C.slate, italic=True)
        # "Your Job"
        add_textbox(slide, x + 0.15, y + 0.7, card_w - 0.3, 0.2,
                    text="Your Job / Tu Trabajo", font_size=9, color=C.amber, bold=True)
        # Jobs (bilingual)
        jobs = role.get("jobs", [])
        jobs_es = role.get("jobs_es", [])
        runs = []
        for j, job in enumerate(jobs):
            runs.append({"text": job, "breakLine": j > 0, "font_size": 9, "color": C.white})
            if j < len(jobs_es):
                runs.append({"text": jobs_es[j], "breakLine": True, "font_size": 8, "color": C.slate, "italic": True})
        add_rich_textbox(slide, x + 0.15, y + 0.95, card_w - 0.3, card_h - 1.1, runs)

    return slide


# ─── STANDARDS-AWARE SLIDES ──────────────────────────────────

def build_vocabulary_cards(prs, data):
    """VOCABULARY_CARDS: Key terms from MCA-III standards for the lesson topic"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.surface)

    add_accent_bar(slide, 0.3, 0.2, 0.45, C.teal)
    add_textbox(slide, 0.5, 0.2, 6, 0.45,
                text="Key Vocabulary", font_size=26,
                font_name=FONT_TITLE, color=C.navy)

    standard_tag = data.get("standard_tag", "")
    if standard_tag:
        add_textbox(slide, 6.5, 0.25, 3.2, 0.35,
                    text=f"MCA-III: {standard_tag}", font_size=10,
                    color=C.slate, italic=True, align=PP_ALIGN.RIGHT)

    terms = data.get("terms", [])
    # Layout: 2-column grid of vocab cards
    cols = 2
    card_w = 4.55
    card_h = 1.1
    gap_x = 0.3
    gap_y = 0.15
    start_x = 0.3
    start_y = 0.85

    for i, term_data in enumerate(terms[:8]):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        # Card background
        add_shape(slide, x, y, card_w, card_h, fill_color=C.white)
        # Accent bar on left
        accent_color = [C.teal, C.amber, C.blue, C.task_mgr][i % 4]
        add_accent_bar(slide, x, y + 0.1, card_h - 0.2, accent_color)

        # Term name
        add_textbox(slide, x + 0.2, y + 0.08, card_w - 0.4, 0.3,
                    text=term_data.get("term", ""),
                    font_size=14, font_name=FONT_HEADING, bold=True, color=C.navy)
        # Definition
        add_textbox(slide, x + 0.2, y + 0.4, card_w - 0.4, 0.6,
                    text=term_data.get("definition", ""),
                    font_size=10, color=C.charcoal)

    return slide


def build_standards_tag(prs, data):
    """STANDARDS_TAG: Small footer overlay showing aligned MCA-III benchmarks.
    Can be added to any slide via post-processing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.surface)

    # Main content placeholder
    add_textbox(slide, 0.5, 0.3, 9, 4.0,
                text=data.get("content", ""),
                font_size=16, color=C.charcoal)

    # Standards tag footer
    benchmarks = data.get("benchmarks", [])
    if benchmarks:
        tag_text = "MCA-III: " + " | ".join(benchmarks[:4])
        add_shape(slide, 0, 5.25, 10, 0.375, fill_color=C.navy)
        add_textbox(slide, 0.3, 5.27, 9.4, 0.3,
                    text=tag_text, font_size=8,
                    color=C.amber, italic=True)

    return slide


def add_standards_footer(slide, benchmarks):
    """Add a MCA-III standards footer bar to any existing slide.
    Call this AFTER building the slide to tag it with aligned standards."""
    if not benchmarks:
        return
    tag_text = "MCA-III: " + " | ".join(benchmarks[:4])
    add_shape(slide, 0, 5.25, 10, 0.375, fill_color=C.navy)
    add_textbox(slide, 0.3, 5.27, 9.4, 0.3,
                text=tag_text, font_size=8,
                color=C.amber, italic=True)


# ─── LESSON DIGESTER: NEW SLIDE BUILDERS ─────────────────────

def build_exit_ticket(prs, data):
    """EXIT_TICKET: 2-3 quick formative assessment problems.
    data keys:
        title (str), problems (list of str), standard_tag (str),
        sentence_frame (str, optional)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.darkBg)

    # Header
    add_shape(slide, 0.3, 0.2, 2.5, 0.55, fill_color=C.amber)
    add_textbox(slide, 0.4, 0.25, 2.3, 0.45,
                text="EXIT TICKET", font_size=22,
                font_name=FONT_TITLE, color=C.white,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    title = data.get("title", "Show What You Know")
    add_textbox(slide, 3.2, 0.25, 6.5, 0.45,
                text=title, font_size=18,
                font_name=FONT_HEADING, color=C.white, bold=True)

    # Problem cards
    problems = data.get("problems", [])
    card_w = 2.85
    gap = 0.25
    start_x = 0.3
    card_y = 1.0
    card_h = 3.2

    for i, prob in enumerate(problems[:3]):
        x = start_x + i * (card_w + gap)
        # Card
        add_shape(slide, x, card_y, card_w, card_h, fill_color=RGBColor(0x2D, 0x37, 0x48))
        # Number badge
        badge_color = [C.teal, C.amber, C.blue][i % 3]
        add_shape(slide, x + 0.15, card_y + 0.15, 0.4, 0.4, fill_color=badge_color)
        add_textbox(slide, x + 0.15, card_y + 0.17, 0.4, 0.35,
                    text=str(i + 1), font_size=18, font_name=FONT_TITLE,
                    color=C.white, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        # Problem text
        add_textbox(slide, x + 0.15, card_y + 0.7, card_w - 0.3, card_h - 0.9,
                    text=prob, font_size=13, color=RGBColor(0xCB, 0xD5, 0xE0))

    # Sentence frame at bottom (optional)
    frame = data.get("sentence_frame")
    if frame:
        add_shape(slide, 0.3, 4.5, 9.4, 0.5, fill_color=RGBColor(0x2D, 0x37, 0x48))
        add_textbox(slide, 0.5, 4.52, 9, 0.45,
                    text=f'"{frame}"', font_size=11,
                    color=C.amber, italic=True, align=PP_ALIGN.CENTER)

    # Standards tag
    tag = data.get("standard_tag", "")
    if tag:
        add_textbox(slide, 6, 5.15, 3.7, 0.3,
                    text=f"MCA-III: {tag}", font_size=8,
                    color=C.slate, italic=True, align=PP_ALIGN.RIGHT)

    return slide


def build_warmup_review(prs, data):
    """WARMUP_REVIEW: Day 2+ opener reviewing yesterday's key concepts.
    data keys:
        title (str), left_title (str), left_items (list of dicts with text, highlight),
        right_title (str), right_items (list of dicts with text, highlight),
        bottom_prompt (str, optional)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.surface)

    add_accent_bar(slide, 0.3, 0.2, 0.45, C.teal)
    add_textbox(slide, 0.5, 0.2, 8, 0.45,
                text=data.get("title", "Yesterday's Key Ideas"),
                font_size=26, font_name=FONT_TITLE, color=C.navy)

    # Left card
    add_shape(slide, 0.3, 0.9, 4.5, 3.3, fill_color=C.white)
    add_textbox(slide, 0.5, 1.0, 4.1, 0.35,
                text=data.get("left_title", "What We Learned"),
                font_size=16, font_name=FONT_HEADING, bold=True, color=C.navy)

    left_items = data.get("left_items", [])
    runs = []
    for j, item in enumerate(left_items):
        text = item.get("text", "")
        highlight = item.get("highlight", "")
        if j > 0:
            runs.append({"text": "", "breakLine": True, "font_size": 6, "color": C.charcoal})
        runs.append({"text": text, "breakLine": j > 0, "font_size": 13, "color": C.charcoal})
        if highlight:
            runs.append({"text": f"\n{highlight}", "breakLine": True, "font_size": 15,
                         "color": C.amber, "bold": True})
    if runs:
        add_rich_textbox(slide, 0.5, 1.5, 4.1, 2.5, runs)

    # Right card
    add_shape(slide, 5.1, 0.9, 4.6, 3.3, fill_color=C.white)
    add_textbox(slide, 5.3, 1.0, 4.2, 0.35,
                text=data.get("right_title", "Key Takeaway"),
                font_size=16, font_name=FONT_HEADING, bold=True, color=C.navy)

    right_items = data.get("right_items", [])
    r_runs = []
    for j, item in enumerate(right_items):
        text = item.get("text", "")
        highlight = item.get("highlight", "")
        if j > 0:
            r_runs.append({"text": "", "breakLine": True, "font_size": 6, "color": C.charcoal})
        r_runs.append({"text": text, "breakLine": j > 0, "font_size": 13, "color": C.charcoal})
        if highlight:
            r_runs.append({"text": f"\n{highlight}", "breakLine": True, "font_size": 15,
                           "color": C.amber, "bold": True})
    if r_runs:
        add_rich_textbox(slide, 5.3, 1.5, 4.2, 2.5, r_runs)

    # Bottom discussion prompt
    prompt = data.get("bottom_prompt")
    if prompt:
        add_shape(slide, 0.3, 4.5, 9.4, 0.8, fill_color=C.lightBlue)
        add_rich_textbox(slide, 0.5, 4.55, 9, 0.7, [
            {"text": "Discuss: ", "bold": True, "font_size": 13, "color": C.navy},
            {"text": prompt, "font_size": 13, "color": C.charcoal},
        ])

    return slide


def build_would_you_rather(prs, data):
    """WOULD_YOU_RATHER: Two-option comparison slide for mathematical reasoning.
    data keys:
        title (str), option_a (dict: label, detail, math, result),
        option_b (dict: same), reveal (str), standard_tag (str)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.white)

    add_textbox(slide, 0.5, 0.15, 9, 0.55,
                text=data.get("title", "Would You Rather...?"),
                font_size=34, font_name=FONT_TITLE,
                color=C.navy, align=PP_ALIGN.CENTER)

    # Option A
    a = data.get("option_a", {})
    add_shape(slide, 0.4, 0.9, 4.3, 3.3, fill_color=C.surface)
    add_textbox(slide, 0.6, 1.0, 3.9, 0.35,
                text=a.get("label", "Option A"), font_size=16,
                font_name=FONT_HEADING, bold=True, color=C.navy, align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.6, 1.4, 3.9, 0.25,
                text=a.get("detail", ""), font_size=11,
                color=C.slate, align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.6, 1.8, 3.9, 1.8,
                text=a.get("math", ""), font_size=14,
                color=C.charcoal, align=PP_ALIGN.CENTER)
    result_a = a.get("result", "")
    if result_a:
        add_textbox(slide, 0.6, 3.6, 3.9, 0.4,
                    text=result_a, font_size=16,
                    font_name=FONT_HEADING, bold=True,
                    color=C.teal, align=PP_ALIGN.CENTER)

    # OR badge
    add_shape(slide, 4.4, 2.0, 1.2, 0.7, fill_color=C.amber)
    add_textbox(slide, 4.4, 2.02, 1.2, 0.65,
                text="OR", font_size=24,
                font_name=FONT_TITLE, color=C.white,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # Option B
    b = data.get("option_b", {})
    add_shape(slide, 5.3, 0.9, 4.3, 3.3, fill_color=C.surface)
    add_textbox(slide, 5.5, 1.0, 3.9, 0.35,
                text=b.get("label", "Option B"), font_size=16,
                font_name=FONT_HEADING, bold=True, color=C.navy, align=PP_ALIGN.CENTER)
    add_textbox(slide, 5.5, 1.4, 3.9, 0.25,
                text=b.get("detail", ""), font_size=11,
                color=C.slate, align=PP_ALIGN.CENTER)
    add_textbox(slide, 5.5, 1.8, 3.9, 1.8,
                text=b.get("math", ""), font_size=14,
                color=C.charcoal, align=PP_ALIGN.CENTER)
    result_b = b.get("result", "")
    if result_b:
        add_textbox(slide, 5.5, 3.6, 3.9, 0.4,
                    text=result_b, font_size=16,
                    font_name=FONT_HEADING, bold=True,
                    color=C.teal, align=PP_ALIGN.CENTER)

    # Reveal bar
    reveal = data.get("reveal")
    if reveal:
        add_shape(slide, 0.5, 4.5, 9, 0.6, fill_color=C.lightBlue)
        add_textbox(slide, 0.7, 4.52, 8.6, 0.55,
                    text=reveal, font_size=12,
                    color=C.navy, italic=True, align=PP_ALIGN.CENTER)

    return slide


def build_notice_wonder(prs, data):
    """NOTICE_WONDER: 'What Do You Notice? What Do You Wonder?' — classic CPM routine.
    data keys:
        title (str), image (str, optional path), context (str),
        notice_starter (str), wonder_starter (str),
        role_prompts (list of (role, prompt) tuples, optional)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.surface)

    # Header
    add_accent_bar(slide, 0.3, 0.2, 0.45, C.amber)
    add_textbox(slide, 0.5, 0.2, 9, 0.45,
                text=data.get("title", "What Do You Notice?  What Do You Wonder?"),
                font_size=24, font_name=FONT_TITLE, color=C.navy)

    # Context / stimulus area (left or top)
    context = data.get("context", "")
    img = data.get("image")

    if img and os.path.exists(img):
        # Image-based: image left, prompts right
        slide.shapes.add_picture(img, Inches(0.4), Inches(0.9), Inches(4.5), Inches(3.5))
        prompt_x = 5.2
        prompt_w = 4.5
    else:
        # Text-based: context card left, prompts right
        add_shape(slide, 0.3, 0.9, 4.5, 3.5, fill_color=C.white)
        add_textbox(slide, 0.5, 1.0, 4.1, 3.2,
                    text=context, font_size=14, color=C.charcoal)
        prompt_x = 5.2
        prompt_w = 4.5

    # NOTICE card
    add_shape(slide, prompt_x, 0.9, prompt_w, 1.5, fill_color=C.white)
    add_shape(slide, prompt_x, 0.9, prompt_w, 0.06, fill_color=C.teal)
    add_textbox(slide, prompt_x + 0.15, 1.05, prompt_w - 0.3, 0.3,
                text="I NOTICE...", font_size=16,
                font_name=FONT_HEADING, bold=True, color=C.teal)
    notice_starter = data.get("notice_starter", "I notice that _____ because _____.")
    add_textbox(slide, prompt_x + 0.15, 1.4, prompt_w - 0.3, 0.8,
                text=f'"{notice_starter}"', font_size=12,
                color=C.charcoal, italic=True)

    # WONDER card
    add_shape(slide, prompt_x, 2.6, prompt_w, 1.5, fill_color=C.white)
    add_shape(slide, prompt_x, 2.6, prompt_w, 0.06, fill_color=C.amber)
    add_textbox(slide, prompt_x + 0.15, 2.75, prompt_w - 0.3, 0.3,
                text="I WONDER...", font_size=16,
                font_name=FONT_HEADING, bold=True, color=C.amber)
    wonder_starter = data.get("wonder_starter", "I wonder if/why/how _____?")
    add_textbox(slide, prompt_x + 0.15, 3.1, prompt_w - 0.3, 0.8,
                text=f'"{wonder_starter}"', font_size=12,
                color=C.charcoal, italic=True)

    # Role prompts at bottom (optional)
    role_strip = data.get("role_prompts")
    if role_strip:
        strip_y = 4.6
        add_shape(slide, 0.3, strip_y, 9.4, 0.7, fill_color=C.white)
        role_colors = [C.facilitator, C.resource_mgr, C.task_mgr, C.recorder]
        for i, (role, prompt) in enumerate(role_strip[:4]):
            rx = 0.4 + i * 2.35
            add_shape(slide, rx, strip_y + 0.08, 0.1, 0.1, fill_color=role_colors[i])
            add_textbox(slide, rx + 0.15, strip_y + 0.03, 2.0, 0.18,
                        text=role, font_size=7, bold=True, color=role_colors[i])
            add_textbox(slide, rx + 0.15, strip_y + 0.22, 2.0, 0.4,
                        text=prompt, font_size=7, color=C.slate, italic=True)

    return slide


def build_learning_log(prs, data):
    """LEARNING_LOG: Reflection / journal prompt slide.
    data keys:
        title (str), prompts (list of str), vocab_check (list of str, optional),
        metacognitive (str, optional — "Rate yourself 1-4" type prompt)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.surface)

    # Header
    add_shape(slide, 0, 0, 10, 0.8, fill_color=C.navy)
    add_textbox(slide, 0.5, 0.1, 9, 0.6,
                text=data.get("title", "Learning Log"),
                font_size=30, font_name=FONT_TITLE,
                color=C.white, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # Reflection prompts (left)
    prompts = data.get("prompts", [])
    add_shape(slide, 0.3, 1.1, 5.5, 3.2, fill_color=C.white)
    add_accent_bar(slide, 0.3, 1.2, 0.35, C.teal)
    add_textbox(slide, 0.5, 1.15, 5.1, 0.35,
                text="Reflect & Write", font_size=16,
                font_name=FONT_HEADING, bold=True, color=C.navy)

    runs = []
    for i, p in enumerate(prompts[:5]):
        if i > 0:
            runs.append({"text": "", "breakLine": True, "font_size": 6, "color": C.charcoal})
        runs.append({"text": f"• {p}", "breakLine": i > 0, "font_size": 13, "color": C.charcoal})
    if runs:
        add_rich_textbox(slide, 0.5, 1.65, 5.1, 2.4, runs)

    # Vocabulary check (right top, optional)
    vocab = data.get("vocab_check", [])
    if vocab:
        add_shape(slide, 6.1, 1.1, 3.6, 1.6, fill_color=C.white)
        add_shape(slide, 6.1, 1.1, 3.6, 0.06, fill_color=C.amber)
        add_textbox(slide, 6.25, 1.2, 3.3, 0.3,
                    text="Vocab Check", font_size=14,
                    font_name=FONT_HEADING, bold=True, color=C.amber)
        add_textbox(slide, 6.25, 1.55, 3.3, 1.0,
                    text="Can you use these words in a sentence?\n" + ", ".join(vocab[:6]),
                    font_size=11, color=C.charcoal)

    # Metacognitive self-rating (right bottom, optional)
    meta = data.get("metacognitive")
    if meta:
        meta_y = 2.9 if vocab else 1.1
        meta_h = 1.4 if vocab else 3.2
        add_shape(slide, 6.1, meta_y, 3.6, meta_h, fill_color=C.white)
        add_shape(slide, 6.1, meta_y, 3.6, 0.06, fill_color=C.blue)
        add_textbox(slide, 6.25, meta_y + 0.15, 3.3, 0.3,
                    text="Rate Yourself", font_size=14,
                    font_name=FONT_HEADING, bold=True, color=C.blue)
        levels = [
            ("Level 1", "I need help, I do not get this"),
            ("Level 2", "I am beginning to make connections"),
            ("Level 3", "I get it with a few small mistakes"),
            ("Level 4", "I could teach someone"),
        ]
        level_runs = []
        for lbl, desc in levels:
            level_runs.append({"text": f"{lbl}: ", "breakLine": len(level_runs) > 0,
                               "font_size": 10, "color": C.navy, "bold": True})
            level_runs.append({"text": desc, "font_size": 10, "color": C.charcoal})
        add_rich_textbox(slide, 6.25, meta_y + 0.5, 3.3, meta_h - 0.7, level_runs)
    elif not vocab:
        # If neither vocab nor metacognitive, fill right side with a tip
        add_shape(slide, 6.1, 1.1, 3.6, 3.2, fill_color=C.lightAmber)
        add_textbox(slide, 6.3, 1.3, 3.2, 2.8,
                    text="Tip: Write at least 3 sentences.\n\nUse vocabulary from today's lesson.\n\nInclude an example or diagram if you can.",
                    font_size=12, color=C.amber)

    return slide


def build_timer(prs, data):
    """TIMER: Visual countdown / work time slide.
    data keys:
        minutes (int), label (str), task (str),
        role_reminders (list of (role, reminder) tuples, optional)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.darkBg)

    minutes = data.get("minutes", 10)
    label = data.get("label", "WORK TIME")

    # Big timer circle (simulated with shapes)
    cx, cy = 5.0, 2.3
    # Outer ring
    add_shape(slide, cx - 1.5, cy - 1.5, 3.0, 3.0, fill_color=RGBColor(0x2D, 0x37, 0x48))
    # Inner circle
    add_shape(slide, cx - 1.2, cy - 1.2, 2.4, 2.4, fill_color=C.darkBg)
    # Time display
    add_textbox(slide, cx - 1.2, cy - 0.6, 2.4, 0.9,
                text=f"{minutes}:00", font_size=44,
                font_name=FONT_TITLE, color=C.amber,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, cx - 1.2, cy + 0.3, 2.4, 0.4,
                text="MINUTES", font_size=12,
                color=C.slate, align=PP_ALIGN.CENTER)

    # Label
    add_textbox(slide, 0.5, 0.3, 9, 0.55,
                text=label, font_size=28,
                font_name=FONT_TITLE, color=C.white, align=PP_ALIGN.CENTER)

    # Task description
    task = data.get("task", "")
    if task:
        add_textbox(slide, 1.5, 0.9, 7, 0.4,
                    text=task, font_size=14,
                    color=C.slate, italic=True, align=PP_ALIGN.CENTER)

    # Role reminders at bottom (optional)
    reminders = data.get("role_reminders")
    if reminders:
        strip_y = 4.3
        role_colors = [C.facilitator, C.resource_mgr, C.task_mgr, C.recorder]
        for i, (role, reminder) in enumerate(reminders[:4]):
            rx = 0.4 + i * 2.4
            add_shape(slide, rx, strip_y, 2.15, 0.9, fill_color=RGBColor(0x2D, 0x37, 0x48))
            add_shape(slide, rx, strip_y, 2.15, 0.04, fill_color=role_colors[i])
            add_textbox(slide, rx + 0.1, strip_y + 0.1, 1.95, 0.2,
                        text=role, font_size=9, bold=True, color=role_colors[i])
            add_textbox(slide, rx + 0.1, strip_y + 0.35, 1.95, 0.45,
                        text=reminder, font_size=8, color=RGBColor(0xCB, 0xD5, 0xE0), italic=True)

    return slide


# ─── MAIN: BUILD DEMO DECK ──────────────────────────────────

def build_demo():
    """Build a complete demo lesson deck."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 1. INTRO AGENDA
    build_intro_agenda(prs, {
        "need": "Notebook, Pencil, Chromebook",
        "do_text": 'Label Your Notebook:\nUnit 4: 7.1.1\n"Distance Rate and Time"',
        "next": "Focus Skill:\n1. Draw a line diagram\n2. Solve for the missing value",
        "focus_title": "Percent Problems",
        "focus_problems": [
            {"title": "Problem 1: Drone Battery", "text": "A delivery drone's battery is at 30%. If it has 50 minutes of flight time remaining, what is the total flight time in minutes?"},
            {"title": "Problem 2: Dinner Tip", "text": "Tristan went out for dinner and the bill was $70. He wanted to leave an 18% tip. What was the total bill?"},
        ],
        "focus_instructions": ["Draw a line diagram", "Solve for the missing value"],
    })

    # 2. CHECK-IN
    build_checkin(prs, {})

    # 3. BREATHING
    build_breathing(prs, {})

    # 4. LEARNING TARGET
    build_learning_target(prs, {
        "targets": [
            {"verb": "create", "rest": "tables, graphs, and rules relating distance to time for objects traveling at a constant rate."},
            {"verb": "determine", "rest": "the relationships between distance, rate, and time."},
        ],
        "success_criteria": [
            "I can gather and record data about a train",
            "I can create a table and graph to show the relationship between distance and time",
            "I can justify my conclusions with tables, graphs and rules",
        ],
        "discuss_prompt": "How is internet speed measured?",
        "discuss_frame": "Internet SPEED is measured with ...",
    })

    # 5. STORY: NARRATIVE
    build_story_narrative(prs, {
        "title": "George's Goodbye",
        "text": '"George stood on the train platform and waved goodbye to his sister as she left for summer camp."\n\nGeorge just watched his sister board the express train. He\'s heading to his car when his CityRail App sends a priority notification:\n\nThe train has malfunctioned and is stuck between stations. With his car and a smartphone, George decides to locate the train himself.',
        "discuss": "How far has the train traveled? What information would help to figure it out?",
    })

    # 6. STORY: DRAMATIC
    build_story_dramatic(prs, {
        "title": "Emergency Broadcast",
        "subtitle": "ALERT\nLINE 7 POWER FAILURE",
        "quote": '"Later, as he was getting in his car to drive home, he saw a light flashing inside the train station and heard an announcement that his sister\'s train had malfunctioned and was stuck on the tracks."',
        "details": [
            "Status: Immobilized",
            "Location: Between Station C and D",
            "Time of Incident: 15 minutes post-departure",
            "No passenger communication available",
        ],
        "role_prompt": 'Resource Manager: What specific notifications are available in our "Announcement" data?',
    })

    # 7. STORY: MISSION (with embedded role prompts)
    build_story_mission(prs, {
        "title": "The Rescue Mission",
        "quote": '"George had a map of the train\'s route and decided to drive to where the train was stuck and pick up his sister."',
        "task_manager_prompt": "How much time do we have before traffic conditions change the rescue window?",
        "role_prompts": [
            {"role": "Facilitator", "prompt": "What information would help George figure out exactly where the train is?"},
            {"role": "Resource Manager", "prompt": "What missing variables (speed, delays, mile markers) might matter?"},
            {"role": "Task Manager", "prompt": "How can we ensure every team member agrees on what he needs to know?"},
            {"role": "Recorder Reporter", "prompt": "Let me write this down... (name) said ______"},
        ],
    })

    # 8. GROUP DISCUSSION (guided conversation)
    build_group_discussion(prs, {
        "title": "Finding the Train",
        "question": "How can George figure out where the train is? What information would help him?",
        "guided_prompts": [
            {"role": "Facilitator", "prompt": "Let's start. Who wants to read the problem? What do we know so far?", "color": C.facilitator},
            {"role": "Resource Manager", "prompt": "What are we supposed to figure out? Do we need any tools?", "color": C.resource_mgr},
            {"role": "Task Manager", "prompt": "What should we try first? Are we ready to move on?", "color": C.task_mgr},
            {"role": "Recorder Reporter", "prompt": "What are we going to say to the class? Let me write it down.", "color": C.recorder},
        ],
        "sentence_frame": "We think the train is at _____ because _____",
    })

    # 9. ACTIVITY LAUNCH
    build_activity_launch(prs, {
        "ref": "Unit 4 - 7.1.1: 7-2 to 7-6",
        "roles": [["Timer", "Marker"], ["Toy Starter", "Measurer"]],
        "needs": ["Chromebook", "Calculator", "Graph paper"],
    })

    # 10. PROBLEM with embedded roles
    build_problem(prs, {
        "number": "7-2",
        "text": "George needs to figure out where the train stopped.\n\nHe knows the train departed at 2:15 PM and was traveling at a constant speed of 60 mph.\n\nThe train stopped 15 minutes after departure.\n\nHow far from the station is the train?",
        "role_prompts": [
            ("Facilitator", "Read the problem. Ask: what do we know?"),
            ("Resource Mgr", "What formula connects D, R, and T?"),
            ("Task Mgr", "Let's set up the equation first."),
            ("Recorder", "Write our answer and method."),
        ],
    })

    # 11. CLOSURE
    build_closure(prs, {
        "prompt": '"Discuss: What are examples of proportional relationships in real life?"',
    })

    # 12. HOMEWORK
    build_homework(prs, {"ref": "7.1.1: 7-10 to 7-14"})

    # Save
    out = "demo-lesson.pptx"
    prs.save(out)
    print(f"Created {out} ({len(prs.slides)} slides)")
    return out


if __name__ == "__main__":
    build_demo()
